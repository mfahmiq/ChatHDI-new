"""
PPTX Generation Service for ChatHDI
Creates PowerPoint presentations from AI-generated content
"""

import os
import io
import base64
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# HDI Brand Colors
HDI_GREEN = RGBColor(16, 185, 129)  # Emerald
HDI_CYAN = RGBColor(6, 182, 212)    # Cyan
HDI_DARK = RGBColor(23, 23, 23)     # Dark bg


class PPTXService:
    """Service for generating PowerPoint presentations"""
    
    def __init__(self):
        pass
    
    def detect_pptx_request(self, message: str) -> bool:
        """Detect if user is requesting PPTX generation"""
        message_lower = message.lower()
        keywords = [
            'buatkan ppt', 'buat ppt', 'buatkan presentasi', 'buat presentasi',
            'generate ppt', 'create ppt', 'generate presentation', 'create presentation',
            'slide presentasi', 'buat slide', 'powerpoint', 'pptx'
        ]
        return any(keyword in message_lower for keyword in keywords)
    
    def create_presentation(self, title: str, slides_content: List[Dict]) -> bytes:
        """
        Create a PowerPoint presentation
        
        Args:
            title: Presentation title
            slides_content: List of slide data with 'title', 'content', 'type'
            
        Returns:
            Bytes of the PPTX file
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        self._add_title_slide(prs, title)
        
        # Content Slides
        for slide_data in slides_content:
            slide_type = slide_data.get('type', 'content')
            if slide_type == 'section':
                self._add_section_slide(prs, slide_data.get('title', ''))
            else:
                self._add_content_slide(
                    prs, 
                    slide_data.get('title', ''),
                    slide_data.get('content', [])
                )
        
        # Thank You Slide
        self._add_thank_you_slide(prs)
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _add_title_slide(self, prs: Presentation, title: str):
        """Add title slide with HDI branding"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(33, 33, 33)
        background.line.fill.background()
        
        # Logo placeholder (H letter)
        logo_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5)
        )
        logo_shape.fill.solid()
        logo_shape.fill.fore_color.rgb = HDI_GREEN
        logo_shape.line.fill.background()
        
        # H text in logo
        logo_text = slide.shapes.add_textbox(Inches(6.35), Inches(1.7), Inches(0.6), Inches(1.1))
        tf = logo_text.text_frame
        p = tf.paragraphs[0]
        p.text = "H"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.5))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "ChatHDI - Hydrogen Development Indonesia"
        p.font.size = Pt(20)
        p.font.color.rgb = HDI_GREEN
        p.alignment = PP_ALIGN.CENTER
    
    def _add_section_slide(self, prs: Presentation, title: str):
        """Add section divider slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(23, 23, 23)
        background.line.fill.background()
        
        # Accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(3.25), Inches(13.333), Inches(1)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = HDI_GREEN
        accent.line.fill.background()
        
        # Section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.35), Inches(11.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, prs: Presentation, title: str, content: List[str]):
        """Add content slide with bullet points"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(33, 33, 33)
        background.line.fill.background()
        
        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(23, 23, 23)
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content bullets
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(229, 231, 235)
            p.space_after = Pt(12)
        
        # Footer line
        footer_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7), Inches(12.333), Inches(0.02)
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = HDI_GREEN
        footer_line.line.fill.background()
    
    def _add_thank_you_slide(self, prs: Presentation):
        """Add thank you/closing slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(23, 23, 23)
        background.line.fill.background()
        
        # Thank you text
        ty_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
        tf = ty_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Terima Kasih"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = HDI_GREEN
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11.333), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Ada pertanyaan?"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(156, 163, 175)
        p.alignment = PP_ALIGN.CENTER
        
        # Contact
        contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1))
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Dibuat dengan ChatHDI"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(107, 114, 128)
        p.alignment = PP_ALIGN.CENTER
    
    async def generate_from_topic(self, topic: str, ai_service) -> Dict:
        """
        Generate a complete presentation from a topic using AI
        
        Args:
            topic: The topic for the presentation
            ai_service: AI service to generate content
            
        Returns:
            Dict with 'success', 'pptx_base64', 'filename', 'error'
        """
        try:
            # Ask AI to generate slide content
            prompt = f"""Buatkan outline presentasi tentang: {topic}

Berikan output dalam format berikut:
JUDUL: [Judul presentasi]

SLIDE 1:
Judul: [Judul slide]
- [Poin 1]
- [Poin 2]
- [Poin 3]

SLIDE 2:
Judul: [Judul slide]
- [Poin 1]
- [Poin 2]
- [Poin 3]

(Lanjutkan sampai 5-7 slide)

Fokus pada konten yang informatif dan terstruktur."""

            messages = [{"role": "user", "content": prompt}]
            response = await ai_service.chat(messages, "hdi-4")
            
            # Parse the response to extract slides
            slides_content = self._parse_ai_response(response)
            
            if not slides_content['slides']:
                return {
                    'success': False,
                    'pptx_base64': None,
                    'filename': None,
                    'error': 'Gagal mengekstrak konten slide dari AI'
                }
            
            # Create the presentation
            pptx_bytes = self.create_presentation(
                slides_content['title'],
                slides_content['slides']
            )
            
            # Convert to base64
            pptx_b64 = base64.b64encode(pptx_bytes).decode('utf-8')
            
            # Generate filename
            safe_topic = "".join(c for c in topic[:30] if c.isalnum() or c in (' ', '-', '_')).strip()
            filename = f"ChatHDI_{safe_topic}.pptx"
            
            return {
                'success': True,
                'pptx_base64': pptx_b64,
                'filename': filename,
                'slides_count': len(slides_content['slides']) + 2,  # +2 for title and thank you
                'error': None
            }
            
        except Exception as e:
            print(f"PPTX generation error: {e}")
            return {
                'success': False,
                'pptx_base64': None,
                'filename': None,
                'error': str(e)
            }
    
    def _parse_ai_response(self, response: str) -> Dict:
        """Parse AI response to extract slides content"""
        result = {
            'title': 'Presentasi ChatHDI',
            'slides': []
        }
        
        lines = response.strip().split('\n')
        current_slide = None
        
        for line in lines:
            line = line.strip()
            
            # Extract title
            if line.upper().startswith('JUDUL:'):
                result['title'] = line.split(':', 1)[1].strip()
            
            # New slide section
            elif line.upper().startswith('SLIDE'):
                if current_slide and current_slide.get('content'):
                    result['slides'].append(current_slide)
                current_slide = {'title': '', 'content': [], 'type': 'content'}
            
            # Slide title
            elif line.lower().startswith('judul:') and current_slide is not None:
                current_slide['title'] = line.split(':', 1)[1].strip()
            
            # Bullet point
            elif line.startswith('-') and current_slide is not None:
                content = line[1:].strip()
                if content:
                    current_slide['content'].append(content)
            
            # Also check for bullet with •
            elif line.startswith('•') and current_slide is not None:
                content = line[1:].strip()
                if content:
                    current_slide['content'].append(content)
        
        # Add last slide
        if current_slide and current_slide.get('content'):
            result['slides'].append(current_slide)
        
        return result


# Singleton instance
pptx_service = PPTXService()
